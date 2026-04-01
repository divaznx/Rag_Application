"""
app.py  —  RAG Web UI Backend (Flask)
=======================================
Serves the UI and exposes API endpoints for:
  POST /api/upload     — upload files, ingest into knowledge base
  POST /api/query      — ask a question
  GET  /api/status     — list indexed folders + chunk counts
  DELETE /api/folder   — remove a folder from the index
"""

import os, json, hashlib, shutil, tempfile
from pathlib import Path
from datetime import datetime
from flask import Flask, request, jsonify, send_from_directory
from flask_cors import CORS

try:
    from dotenv import load_dotenv
    load_dotenv()
except ImportError:
    pass

# ── import RAG components ──────────────────────────────────────
import chromadb
from embedder import get_embedding_function
from openai import OpenAI

# ── config ────────────────────────────────────────────────────
KB_DIR          = Path("knowledge_base")
PROCESSED_DIR   = KB_DIR / "processed"
REGISTRY_FILE   = KB_DIR / "registry.json"
COLLECTION_NAME = "rag_collection"
UPLOAD_DIR      = KB_DIR / "uploads"
CHUNK_SIZE      = 800
CHUNK_OVERLAP   = 150
DEFAULT_MODEL   = "gpt-4o-mini"
TOP_K           = 5
SIMILARITY_CUTOFF = 1.4

SUPPORTED_EXTS = {
    ".txt", ".md", ".rst", ".log", ".csv", ".json", ".yaml", ".yml",
    ".py", ".js", ".ts", ".html", ".xml", ".css", ".sql",
    ".pdf", ".docx", ".pptx", ".xlsx", ".xls", ".epub",
}

app = Flask(__name__, static_folder="static")
CORS(app)

# ── helpers ───────────────────────────────────────────────────
def get_collection():
    PROCESSED_DIR.mkdir(parents=True, exist_ok=True)
    client = chromadb.PersistentClient(path=str(PROCESSED_DIR))
    ef = get_embedding_function()
    return client.get_or_create_collection(
        name=COLLECTION_NAME,
        embedding_function=ef,
        metadata={"hnsw:space": "cosine"},
    )

def load_registry():
    if REGISTRY_FILE.exists():
        with open(REGISTRY_FILE) as f:
            return json.load(f)
    return {}

def save_registry(reg):
    REGISTRY_FILE.parent.mkdir(parents=True, exist_ok=True)
    with open(REGISTRY_FILE, "w") as f:
        json.dump(reg, f, indent=2)

def extract_text(filepath: Path):
    ext = filepath.suffix.lower()
    if ext not in SUPPORTED_EXTS:
        return None
    try:
        if ext in {".txt",".md",".rst",".log",".csv",".json",".yaml",".yml",
                   ".py",".js",".ts",".html",".xml",".css",".sql"}:
            return filepath.read_text(encoding="utf-8", errors="ignore")
        elif ext == ".pdf":
            from pypdf import PdfReader
            r = PdfReader(str(filepath))
            return "\n\n".join(p.extract_text() or "" for p in r.pages).strip() or None
        elif ext == ".docx":
            import docx as dx
            doc = dx.Document(str(filepath))
            parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
            for t in doc.tables:
                for row in t.rows:
                    for cell in row.cells:
                        if cell.text.strip(): parts.append(cell.text.strip())
            return "\n\n".join(parts) or None
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(filepath))
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                texts = [s.text.strip() for s in slide.shapes if hasattr(s,"text") and s.text.strip()]
                if texts: slides.append(f"[Slide {i}]\n" + "\n".join(texts))
            return "\n\n".join(slides) or None
        elif ext in {".xlsx",".xls"}:
            import openpyxl
            wb = openpyxl.load_workbook(str(filepath), read_only=True, data_only=True)
            sheets = []
            for name in wb.sheetnames:
                rows = ["\t".join(str(c) if c is not None else "" for c in row)
                        for row in wb[name].iter_rows(values_only=True)]
                rows = [r for r in rows if r.strip()]
                if rows: sheets.append(f"[Sheet: {name}]\n" + "\n".join(rows))
            return "\n\n".join(sheets) or None
        elif ext == ".epub":
            import ebooklib
            from ebooklib import epub
            from bs4 import BeautifulSoup
            book = epub.read_epub(str(filepath))
            parts = []
            for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
                text = BeautifulSoup(item.get_content(), "html.parser").get_text("\n").strip()
                if text: parts.append(text)
            return "\n\n".join(parts) or None
    except Exception as e:
        print(f"  [warn] {filepath.name}: {e}")
        return None

def chunk_text(text):
    if not text.strip(): return []
    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
    chunks, current = [], ""
    for para in paragraphs:
        if len(current) + len(para) + 2 <= CHUNK_SIZE:
            current = (current + "\n\n" + para).strip()
        else:
            if current: chunks.append(current)
            if len(para) > CHUNK_SIZE:
                sentences = para.replace(". ",".|").replace("? ","?|").replace("! ","!|").split("|")
                buf = ""
                for sent in sentences:
                    if len(buf) + len(sent) + 1 <= CHUNK_SIZE:
                        buf = (buf + " " + sent).strip()
                    else:
                        if buf: chunks.append(buf)
                        buf = sent
                current = buf
            else:
                current = para
    if current: chunks.append(current)
    result = []
    for i, chunk in enumerate(chunks):
        if i > 0 and CHUNK_OVERLAP > 0:
            chunk = chunks[i-1][-CHUNK_OVERLAP:].rstrip() + "\n" + chunk.lstrip()
        result.append(chunk)
    return result

def ingest_file(filepath: Path, folder_name: str, collection):
    text = extract_text(filepath)
    if not text: return 0
    chunks = chunk_text(text)
    if not chunks: return 0
    doc_ids, texts, metas = [], [], []
    for i, chunk in enumerate(chunks):
        doc_id = hashlib.md5(f"{folder_name}/{filepath.name}/{i}".encode()).hexdigest()
        doc_ids.append(doc_id)
        texts.append(chunk)
        metas.append({"folder": folder_name, "file": filepath.name,
                      "chunk_id": i, "ext": filepath.suffix.lower()})
    for start in range(0, len(texts), 100):
        collection.upsert(ids=doc_ids[start:start+100],
                          documents=texts[start:start+100],
                          metadatas=metas[start:start+100])
    return len(chunks)

# ── API routes ────────────────────────────────────────────────

@app.route("/api/status", methods=["GET"])
def status():
    try:
        collection = get_collection()
        registry   = load_registry()
        return jsonify({
            "total_chunks": collection.count(),
            "folders": [
                {"name": name, **meta}
                for name, meta in registry.items()
            ]
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/upload", methods=["POST"])
def upload():
    files = request.files.getlist("files")
    if not files:
        return jsonify({"error": "No files provided"}), 400

    collection = get_collection()
    registry   = load_registry()
    results    = []
    UPLOAD_DIR.mkdir(parents=True, exist_ok=True)

    # group files under a single "upload session" folder name
    folder_name = request.form.get("folder_name") or \
                  f"upload_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
    folder_path = UPLOAD_DIR / folder_name
    folder_path.mkdir(parents=True, exist_ok=True)

    total_chunks = 0
    for f in files:
        if not f.filename: continue
        ext = Path(f.filename).suffix.lower()
        if ext not in SUPPORTED_EXTS:
            results.append({"file": f.filename, "status": "skipped",
                             "reason": f"Unsupported type {ext}"})
            continue
        dest = folder_path / f.filename
        f.save(str(dest))
        chunks = ingest_file(dest, folder_name, collection)
        total_chunks += chunks
        results.append({"file": f.filename, "status": "ok", "chunks": chunks})

    if total_chunks > 0:
        existing = registry.get(folder_name, {"files": 0, "chunks": 0})
        registry[folder_name] = {
            "hash":       "",
            "chunks":     existing["chunks"] + total_chunks,
            "files":      existing["files"]  + len([r for r in results if r["status"]=="ok"]),
            "indexed_at": datetime.now().isoformat(timespec="seconds"),
        }
        save_registry(registry)

    return jsonify({"folder": folder_name, "results": results,
                    "total_chunks": total_chunks})


@app.route("/api/query", methods=["POST"])
def query():
    data  = request.get_json()
    question = (data or {}).get("question", "").strip()
    folder   = (data or {}).get("folder")
    model    = (data or {}).get("model", DEFAULT_MODEL)
    if not question:
        return jsonify({"error": "No question provided"}), 400

    try:
        collection = get_collection()
        if collection.count() == 0:
            return jsonify({"answer": "The knowledge base is empty. Please upload some documents first.",
                            "sources": []})

        where  = {"folder": folder} if folder else None
        kwargs = dict(query_texts=[question],
                      n_results=min(TOP_K, collection.count()),
                      include=["documents","metadatas","distances"])
        if where: kwargs["where"] = where
        results = collection.query(**kwargs)

        docs      = results["documents"][0]
        metas     = results["metadatas"][0]
        distances = results["distances"][0]

        chunks = [{"text": doc, "folder": m.get("folder","?"),
                   "file": m.get("file","?"), "distance": round(d,4)}
                  for doc, m, d in zip(docs, metas, distances)
                  if d <= SIMILARITY_CUTOFF]

        if not chunks:
            return jsonify({"answer": "I couldn't find relevant information in the knowledge base for your question.",
                            "sources": []})

        context = "\n\n---\n\n".join(
            f"[Source: {c['folder']}/{c['file']}]\n{c['text']}" for c in chunks
        )
        messages = [
            {"role": "system", "content":
             "You are a helpful assistant. Answer strictly based on the provided context. "
             "If the context lacks enough info, say so. Be concise and precise."},
            {"role": "user", "content":
             f"Context:\n{context}\n\nQuestion: {question}"}
        ]

        client   = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
        response = client.chat.completions.create(
            model=model, messages=messages, temperature=0.1, max_tokens=1024
        )
        answer  = response.choices[0].message.content.strip()
        sources = list({f"{c['folder']}/{c['file']}" for c in chunks})
        return jsonify({"answer": answer, "sources": sources,
                        "chunks_used": len(chunks)})

    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/folder/<folder_name>", methods=["DELETE"])
def delete_folder(folder_name):
    try:
        collection = get_collection()
        registry   = load_registry()
        results    = collection.get(where={"folder": folder_name})
        ids        = results.get("ids", [])
        if ids: collection.delete(ids=ids)
        registry.pop(folder_name, None)
        save_registry(registry)
        folder_path = UPLOAD_DIR / folder_name
        if folder_path.exists(): shutil.rmtree(folder_path)
        return jsonify({"deleted": folder_name, "chunks_removed": len(ids)})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/", defaults={"path": ""})
@app.route("/<path:path>")
def serve(path):
    static_dir = Path("static")
    if path and (static_dir / path).exists():
        return send_from_directory("static", path)
    return send_from_directory("static", "index.html")


if __name__ == "__main__":
    Path("static").mkdir(exist_ok=True)
    print("\n  RAG Web UI running at  →  http://localhost:5000\n")
    app.run(debug=True, port=5000)