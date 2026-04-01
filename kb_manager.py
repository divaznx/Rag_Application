"""
kb_manager.py  —  Knowledge Base Manager
=========================================
Ingests any folder into ChromaDB, skipping already-processed ones.

Supported formats:
  Text    : .txt  .md  .rst  .log  .csv  .json  .yaml  .yml
  Code    : .py  .js  .ts  .html  .xml  .css  .sql
  PDF     : .pdf
  Word    : .docx
  PowerPoint: .pptx
  Excel   : .xlsx  .xls
  Ebook   : .epub

Usage:
    python kb_manager.py --ingest ./my_docs
    python kb_manager.py --ingest ./my_docs --force
    python kb_manager.py --status
    python kb_manager.py --delete <folder_name>
"""

import os, sys, json, hashlib, argparse
from pathlib import Path
from datetime import datetime

import chromadb
from embedder import get_embedding_function

try:
    from dotenv import load_dotenv
    load_dotenv()
except ImportError:
    pass

# ─────────────────────────────────────────────
#  CONFIG
# ─────────────────────────────────────────────
KB_DIR          = Path("knowledge_base")
PROCESSED_DIR   = KB_DIR / "processed"
REGISTRY_FILE   = KB_DIR / "registry.json"
COLLECTION_NAME = "rag_collection"

CHUNK_SIZE    = 800
CHUNK_OVERLAP = 150

SUPPORTED_EXTS = {
    # plain text / markup
    ".txt", ".md", ".rst", ".log", ".csv",
    ".json", ".yaml", ".yml",
    # code
    ".py", ".js", ".ts", ".html", ".xml",
    ".css", ".sql",
    # documents
    ".pdf", ".docx", ".pptx", ".xlsx", ".xls",
    # ebook
    ".epub",
}

# ─────────────────────────────────────────────
#  COLOURS
# ─────────────────────────────────────────────
class C:
    RESET="\033[0m"; BOLD="\033[1m"; GREEN="\033[92m"
    YELLOW="\033[93m"; RED="\033[91m"; CYAN="\033[96m"
    BLUE="\033[94m"; DIM="\033[2m"

def banner(text):
    line = "─" * 60
    print(f"\n{C.BOLD}{C.CYAN}{line}\n  {text}\n{line}{C.RESET}\n")

def ok(msg):   print(f"  {C.GREEN}✔  {C.RESET}{msg}")
def warn(msg): print(f"  {C.YELLOW}⚠  {C.RESET}{msg}")
def err(msg):  print(f"  {C.RED}✘  {C.RESET}{msg}")
def info(msg): print(f"  {C.BLUE}ℹ  {C.RESET}{msg}")


# ─────────────────────────────────────────────
#  TEXT EXTRACTION  (per file type)
# ─────────────────────────────────────────────
def extract_text(filepath: Path) -> str | None:
    ext = filepath.suffix.lower()

    if ext not in SUPPORTED_EXTS:
        return None

    try:
        # ── Plain text / code / markup ──────────────────────────
        if ext in {".txt", ".md", ".rst", ".log", ".csv", ".json",
                   ".yaml", ".yml", ".py", ".js", ".ts", ".html",
                   ".xml", ".css", ".sql"}:
            return filepath.read_text(encoding="utf-8", errors="ignore")

        # ── PDF ─────────────────────────────────────────────────
        elif ext == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(str(filepath))
            pages  = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    pages.append(text.strip())
            return "\n\n".join(pages) if pages else None

        # ── Word (.docx) ─────────────────────────────────────────
        elif ext == ".docx":
            import docx as docx_lib
            doc   = docx_lib.Document(str(filepath))
            paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
            # also grab text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            paras.append(cell.text.strip())
            return "\n\n".join(paras) if paras else None

        # ── PowerPoint (.pptx) ───────────────────────────────────
        elif ext == ".pptx":
            from pptx import Presentation
            prs    = Presentation(str(filepath))
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                if texts:
                    slides.append(f"[Slide {i}]\n" + "\n".join(texts))
            return "\n\n".join(slides) if slides else None

        # ── Excel (.xlsx / .xls) ─────────────────────────────────
        elif ext in {".xlsx", ".xls"}:
            import openpyxl
            wb     = openpyxl.load_workbook(str(filepath), read_only=True, data_only=True)
            sheets = []
            for sheet_name in wb.sheetnames:
                ws   = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    row_text = "\t".join(str(c) if c is not None else "" for c in row)
                    if row_text.strip():
                        rows.append(row_text)
                if rows:
                    sheets.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
            return "\n\n".join(sheets) if sheets else None

        # ── EPUB ─────────────────────────────────────────────────
        elif ext == ".epub":
            import ebooklib
            from ebooklib import epub
            from bs4 import BeautifulSoup
            book   = epub.read_epub(str(filepath))
            parts  = []
            for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
                soup = BeautifulSoup(item.get_content(), "html.parser")
                text = soup.get_text(separator="\n").strip()
                if text:
                    parts.append(text)
            return "\n\n".join(parts) if parts else None

    except Exception as e:
        warn(f"Could not read {filepath.name}: {e}")
        return None


# ─────────────────────────────────────────────
#  CHUNKING
# ─────────────────────────────────────────────
def chunk_text(text: str, chunk_size=CHUNK_SIZE, overlap=CHUNK_OVERLAP) -> list[str]:
    if not text.strip():
        return []

    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
    chunks, current = [], ""

    for para in paragraphs:
        if len(current) + len(para) + 2 <= chunk_size:
            current = (current + "\n\n" + para).strip()
        else:
            if current:
                chunks.append(current)
            if len(para) > chunk_size:
                # split long paragraphs by sentence
                sentences = para.replace(". ",".|").replace("? ","?|").replace("! ","!|").split("|")
                buf = ""
                for sent in sentences:
                    if len(buf) + len(sent) + 1 <= chunk_size:
                        buf = (buf + " " + sent).strip()
                    else:
                        if buf:
                            chunks.append(buf)
                        buf = sent
                current = buf
            else:
                current = para

    if current:
        chunks.append(current)

    # apply overlap
    result = []
    for i, chunk in enumerate(chunks):
        if i > 0 and overlap > 0:
            chunk = chunks[i-1][-overlap:].rstrip() + "\n" + chunk.lstrip()
        result.append(chunk)

    return result


# ─────────────────────────────────────────────
#  FOLDER HASH  (change detection)
# ─────────────────────────────────────────────
def folder_hash(folder: Path) -> str:
    h = hashlib.md5()
    for fp in sorted(folder.rglob("*")):
        if fp.is_file() and fp.suffix.lower() in SUPPORTED_EXTS:
            h.update(str(fp.relative_to(folder)).encode())
            h.update(str(fp.stat().st_size).encode())
    return h.hexdigest()


# ─────────────────────────────────────────────
#  REGISTRY
# ─────────────────────────────────────────────
def load_registry() -> dict:
    if REGISTRY_FILE.exists():
        with open(REGISTRY_FILE) as f:
            return json.load(f)
    return {}

def save_registry(reg: dict):
    REGISTRY_FILE.parent.mkdir(parents=True, exist_ok=True)
    with open(REGISTRY_FILE, "w") as f:
        json.dump(reg, f, indent=2)


# ─────────────────────────────────────────────
#  CHROMADB
# ─────────────────────────────────────────────
def get_collection():
    PROCESSED_DIR.mkdir(parents=True, exist_ok=True)
    client = chromadb.PersistentClient(path=str(PROCESSED_DIR))
    ef     = get_embedding_function()
    return client.get_or_create_collection(
        name=COLLECTION_NAME,
        embedding_function=ef,
        metadata={"hnsw:space": "cosine"},
    )


# ─────────────────────────────────────────────
#  INGEST ONE FOLDER
# ─────────────────────────────────────────────
def ingest_folder(folder: Path, collection, registry: dict, force=False) -> bool:
    folder_name  = folder.name
    current_hash = folder_hash(folder)

    if not force and folder_name in registry:
        if registry[folder_name].get("hash") == current_hash:
            warn(f"'{folder_name}' already indexed and unchanged — skipping.")
            return False
        else:
            info(f"'{folder_name}' has changed — re-indexing.")
            delete_folder_chunks(folder_name, collection, registry, quiet=True)

    files = [fp for fp in folder.rglob("*")
             if fp.is_file() and fp.suffix.lower() in SUPPORTED_EXTS]

    if not files:
        warn(f"'{folder_name}' has no supported files — skipping.")
        return False

    info(f"Ingesting '{folder_name}' ({len(files)} file(s))…")
    doc_ids, texts, metas = [], [], []
    total_chunks = 0

    for fp in files:
        text = extract_text(fp)
        if not text:
            warn(f"  {fp.name} — no text extracted, skipping.")
            continue

        chunks = chunk_text(text)
        if not chunks:
            continue

        for i, chunk in enumerate(chunks):
            doc_id = hashlib.md5(
                f"{folder_name}/{fp.relative_to(folder)}/{i}".encode()
            ).hexdigest()
            doc_ids.append(doc_id)
            texts.append(chunk)
            metas.append({
                "folder":   folder_name,
                "file":     fp.name,
                "chunk_id": i,
                "ext":      fp.suffix.lower(),
            })

        total_chunks += len(chunks)
        ok(f"  {fp.name}  [{fp.suffix}]  →  {len(chunks)} chunk(s)")

    if not texts:
        warn(f"No text extracted from '{folder_name}'.")
        return False

    # upsert in batches of 100
    for start in range(0, len(texts), 100):
        collection.upsert(
            ids=doc_ids[start:start+100],
            documents=texts[start:start+100],
            metadatas=metas[start:start+100],
        )

    registry[folder_name] = {
        "hash":       current_hash,
        "chunks":     total_chunks,
        "files":      len(files),
        "indexed_at": datetime.now().isoformat(timespec="seconds"),
    }
    ok(f"'{folder_name}' done — {total_chunks} chunks indexed.\n")
    return True


# ─────────────────────────────────────────────
#  DELETE FOLDER FROM INDEX
# ─────────────────────────────────────────────
def delete_folder_chunks(folder_name: str, collection, registry: dict, quiet=False):
    results = collection.get(where={"folder": folder_name})
    ids = results.get("ids", [])
    if ids:
        collection.delete(ids=ids)
        if not quiet:
            ok(f"Removed {len(ids)} chunk(s) for '{folder_name}'.")
    registry.pop(folder_name, None)


# ─────────────────────────────────────────────
#  STATUS
# ─────────────────────────────────────────────
def print_status(registry: dict, collection):
    banner("Knowledge Base Status")
    if not registry:
        warn("Nothing indexed yet.")
        return
    total = collection.count()
    print(f"  {C.BOLD}Total chunks in ChromaDB:{C.RESET} {total}\n")
    print(f"  {'Folder':<30} {'Files':>6} {'Chunks':>8}  Indexed At")
    print(f"  {'─'*30} {'─'*6} {'─'*8}  {'─'*19}")
    for name, meta in registry.items():
        print(f"  {C.CYAN}{name:<30}{C.RESET}"
              f"{meta['files']:>6} {meta['chunks']:>8}  "
              f"{C.DIM}{meta['indexed_at']}{C.RESET}")
    print()


# ─────────────────────────────────────────────
#  MAIN
# ─────────────────────────────────────────────
def main():
    parser = argparse.ArgumentParser(description="Knowledge Base Manager")
    group  = parser.add_mutually_exclusive_group(required=True)
    group.add_argument("--ingest", metavar="FOLDER", help="Folder to ingest")
    group.add_argument("--status", action="store_true", help="Show indexed folders")
    group.add_argument("--delete", metavar="FOLDER_NAME", help="Remove folder from index")
    parser.add_argument("--force", action="store_true", help="Re-index even if unchanged")
    args = parser.parse_args()

    banner("RAG — Knowledge Base Manager")
    collection = get_collection()
    registry   = load_registry()

    if args.status:
        print_status(registry, collection)
        return

    if args.delete:
        delete_folder_chunks(args.delete, collection, registry)
        save_registry(registry)
        info(f"'{args.delete}' removed from index.")
        return

    if args.ingest:
        ingest_root = Path(args.ingest).resolve()
        if not ingest_root.exists():
            err(f"Path not found: {ingest_root}")
            sys.exit(1)

        sub_folders = [p for p in ingest_root.iterdir() if p.is_dir()]
        has_direct  = any(
            p.is_file() and p.suffix.lower() in SUPPORTED_EXTS
            for p in ingest_root.iterdir()
        )

        targets = []
        if has_direct:
            targets.append(ingest_root)
        targets.extend(sub_folders)

        if not targets:
            warn("No supported files found in the given path.")
            sys.exit(0)

        ingested = sum(
            ingest_folder(f, collection, registry, force=args.force)
            for f in targets
        )
        save_registry(registry)
        print()
        if ingested:
            ok(f"Done! {ingested} folder(s) ingested.")
        else:
            info("Nothing new to ingest.")
        print_status(registry, collection)


if __name__ == "__main__":
    main()